# -*- coding: utf-8 -*-
"""
기계연구원 샘플 실험 결과 — 재구성 초안
주제: 레이저 초음파(pump-probe) 기반 시편 미세조직(결정립) 비파괴 분석
원본(기계연구원 샘플 실험_260609.pptx) 내용 + 보유 .dat 원본 데이터 분석 결과 반영
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------- 테마 ----------------
NAVY   = RGBColor(0x1F, 0x2A, 0x44)
BLUE   = RGBColor(0x2E, 0x86, 0xDE)
TEAL   = RGBColor(0x12, 0xA5, 0x9B)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF9)
GRAY   = RGBColor(0x55, 0x5A, 0x66)
DGRAY  = RGBColor(0x33, 0x37, 0x40)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT2= RGBColor(0xF5, 0xA6, 0x23)
ORANGE = RGBColor(0xD9, 0x53, 0x1E)

FONT   = "맑은 고딕"
MONO   = "Consolas"

ROOT   = os.path.dirname(os.path.abspath(__file__))
RESDIR = os.path.join(ROOT, "analysis_matlab", "results")
ASSET  = os.path.join(ROOT, "ppt_draft_assets")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_pageno = [0]


def _set(p, text, size, color, bold=False, font=FONT, italic=False):
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic


def rect(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def header(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Pt(4), BLUE)
    tb, tf = textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.95),
                     anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        _set(tf.paragraphs[0], kicker, 12, BLUE, bold=True)
        p1 = tf.add_paragraph(); _set(p1, title, 26, WHITE, bold=True)
    else:
        _set(tf.paragraphs[0], title, 28, WHITE, bold=True)


def footer(slide):
    _pageno[0] += 1
    tb, tf = textbox(slide, Inches(0.5), Inches(7.05), Inches(11), Inches(0.35))
    _set(tf.paragraphs[0],
         "기계연구원 샘플 실험 결과   ·   레이저 초음파 미세조직 분석   ·   26.06.09", 9, GRAY)
    tb2, tf2 = textbox(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _set(tf2.paragraphs[0], str(_pageno[0]), 10, GRAY, bold=True)


def bullets(slide, items, x, y, w, h, size=16, gap=10):
    tb, tf = textbox(slide, x, y, w, h)
    for i, (lvl, txt, *col) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        c = col[0] if col else (DGRAY if lvl == 0 else GRAY)
        bullet = "●  " if lvl == 0 else "–  "
        _set(p, bullet + txt, size if lvl == 0 else size-2, c, bold=(lvl == 0))
        p.level = lvl
        p.space_after = Pt(gap if lvl == 0 else gap-4)
        p.line_spacing = 1.12
    return tb


def card(slide, x, y, w, h, title, body, accent=BLUE, ts=15, bs=12):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, Inches(0.12), h, accent)
    tb, tf = textbox(slide, x + Inches(0.28), y + Inches(0.12),
                     w - Inches(0.45), h - Inches(0.2))
    _set(tf.paragraphs[0], title, ts, NAVY, bold=True)
    for line in body:
        p = tf.add_paragraph(); _set(p, line, bs, GRAY); p.line_spacing = 1.12
        p.space_before = Pt(3)


def caption(slide, x, y, w, text, size=11):
    tb, tf = textbox(slide, x, y, w, Inches(0.35))
    _set(tf.paragraphs[0], text, size, GRAY, italic=True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def pic_fit(slide, path, x, y, w, h):
    """주어진 박스(w,h) 안에 비율 유지하며 가운데 배치."""
    from PIL import Image
    if not os.path.exists(path):
        rect(slide, x, y, w, h, LIGHT)
        return
    iw, ih = Image.open(path).size
    bw, bh = float(w), float(h)
    s = min(bw/iw, bh/ih)
    nw, nh = iw*s, ih*s
    nx = x + (bw-nw)/2; ny = y + (bh-nh)/2
    slide.shapes.add_picture(path, int(nx), int(ny), int(nw), int(nh))


# ======================================================================
# Slide 1 — 표지
# ======================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.35), SW, Pt(3), BLUE)
tb, tf = textbox(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(3))
_set(tf.paragraphs[0], "기계연구원 샘플 실험 결과", 20, BLUE, bold=True)
p = tf.add_paragraph()
_set(p, "레이저 초음파 기반\n시편 미세조직 비파괴 분석", 40, WHITE, bold=True)
p.space_before = Pt(8); p.line_spacing = 1.05
p = tf.add_paragraph()
_set(p, "pump–probe 표면파(SAW)로 결정립 특성 살펴보기", 18, RGBColor(0xC9,0xD6,0xEA))
p.space_before = Pt(16)
tb, tf = textbox(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2))
_set(tf.paragraphs[0], "한양대학교 박성현 교수님 연구실   ·   음파 자성 시편", 15, ACCENT2, bold=True)
p = tf.add_paragraph()
_set(p, "2026. 06. 09", 13, RGBColor(0xAE,0xB9,0xCC)); p.space_before = Pt(8)

# ======================================================================
# Slide 2 — 목차
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "목차", "CONTENTS")
items = [
    ("1", "실험 목적", "미세조직(결정립) 비파괴 분석"),
    ("2", "측정 원리", "왜 레이저 초음파로 결정립을 보는가"),
    ("3", "실험 세팅", "pump–probe 구성 · 레이저 스펙"),
    ("4", "측정 방법", "스캔 방식 · 표면파 전파"),
    ("5", "실험 결과", "취득 신호 · 시편 간 비교"),
    ("6", "결과 분석 & 추후 방향", "신호 품질 이슈 · 개선안"),
]
y = Inches(1.5)
for num, t, d in items:
    rect(s, Inches(0.8), y, Inches(0.62), Inches(0.62), BLUE)
    tb, tf = textbox(s, Inches(0.8), y, Inches(0.62), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set(tf.paragraphs[0], num, 18, WHITE, bold=True)
    tb, tf = textbox(s, Inches(1.65), y, Inches(10.8), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], t + "    —    " + d, 18, DGRAY, bold=True)
    y += Inches(0.88)
footer(s)

# ======================================================================
# Slide 3 — 실험 목적
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "실험 목적", "OBJECTIVE")
bullets(s, [
    (0, "레이저 초음파(LU)로 시편의 미세조직을 비파괴적으로 분석"),
    (1, "결함 검출이 아니라, 동일 공정 시편의 결정립(grain) 특성 재확인이 목표"),
    (0, "결정립 크기가 크게 다른 두 시편을 비교 대상으로 선정"),
], Inches(0.7), Inches(1.45), Inches(12), Inches(1.7), size=17, gap=10)
# 두 시편 카드 + 사진
y0 = Inches(3.1)
card(s, Inches(0.7), y0, Inches(5.9), Inches(3.4),
     "시편 A  (test #5-03)",
     ["", "· 결정립 크기(면적 기준): 64.4 µm  — 미세립",
      "· 시편 크기: 약 6 × 11 mm", "",
      "* Ablation(레이저 파워) 조건 확인 중", "  중앙부에 열손상 발생"], BLUE, ts=16, bs=13)
card(s, Inches(6.85), y0, Inches(5.9), Inches(3.4),
     "시편 B  (test #5-23)",
     ["", "· 결정립 크기(면적 기준): 482 µm  — 조대립",
      "· 시편 크기: 약 6 × 11 mm", "",
      "→ A 대비 약 7.5배 큰 결정립", "  ⇒ 초음파 산란/감쇠 차이 기대"], ORANGE, ts=16, bs=13)
footer(s)

# ======================================================================
# Slide 4 — 측정 원리
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "측정 원리 — 왜 레이저 초음파로 결정립을 보는가", "PRINCIPLE")
bullets(s, [
    (0, "발생 레이저 펄스가 시편 표면을 순간 가열 → 표면파(Rayleigh/SAW) 발생"),
    (0, "3 mm 떨어진 지점을 측정 레이저로 관측 → 표면파의 전파 신호 취득"),
    (0, "결정립 크기가 표면파에 남기는 흔적"),
    (1, "산란·감쇠: 결정립이 클수록 산란 증가 → 고주파 성분이 더 많이 감쇠"),
    (1, "주파수: 고주파 감쇠 ⇒ 스펙트럼 중심주파수 하락 경향"),
    (1, "속도: 미세조직에 따라 표면파 속도(도달시간)도 미세하게 달라질 수 있음"),
], Inches(0.7), Inches(1.5), Inches(7.4), Inches(5), size=16, gap=11)
card(s, Inches(8.4), Inches(1.55), Inches(4.35), Inches(4.7),
     "핵심 아이디어",
     ["", "표면파의", "  · 감쇠", "  · 주파수 성분", "  · 전파 속도", "",
      "를 시편 간 비교하면", "결정립 차이를 비파괴로", "간접 측정할 수 있다.", "",
      "→ 본 실험은 그 가능성을", "   타진하는 1차 측정"], TEAL, ts=17, bs=14)
footer(s)

# ======================================================================
# Slide 5 — 실험 세팅
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "실험 세팅 — pump–probe 구성", "SETUP")
pic_fit(s, os.path.join(ASSET, "photo_setup.png"),
        Inches(0.7), Inches(1.45), Inches(5.2), Inches(4.6))
caption(s, Inches(0.7), Inches(6.05), Inches(5.2), "실제 광학 세팅 (발생/측정 레이저·미러·렌즈)")
# 구성 설명 + 레이저 스펙 표
card(s, Inches(6.2), Inches(1.45), Inches(6.5), Inches(2.05),
     "빔 경로",
     ["발생용 레이저 → 미러 → 오목렌즈 → 시편 표면 (초음파 발생)",
      "측정용 레이저 → 미러 → 시편 표면 (표면파 관측)",
      "펌프–프로브 빔 간격 : 3 mm"], BLUE, ts=15, bs=13)
# 레이저 스펙 표
ty = Inches(3.75); rows = [
    ("항목", "발생용 레이저", "측정용 레이저"),
    ("최대 파워", "3 MW", "– (CW)"),
    ("펄스 폭", "3–5 ns", "–"),
    ("파장", "532 nm", "532 nm"),
    ("빔 직경", "3 mm", "250 µm"),
]
cw = [Inches(2.0), Inches(2.25), Inches(2.25)]; x0 = Inches(6.2); rh = Inches(0.5)
for ri, row in enumerate(rows):
    x = x0
    for ci, val in enumerate(row):
        fill = NAVY if ri == 0 else (LIGHT if ri % 2 else WHITE)
        rect(s, x, ty, cw[ci], rh, fill, line=RGBColor(0xD0,0xD6,0xDE))
        tb, tf = textbox(s, x, ty, cw[ci], rh, anchor=MSO_ANCHOR.MIDDLE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        col = WHITE if ri == 0 else DGRAY
        _set(tf.paragraphs[0], val, 12, col, bold=(ri == 0 or ci == 0))
        x += cw[ci]
    ty += rh
footer(s)

# ======================================================================
# Slide 6 — 측정 방법
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "측정 방법 — 라인 스캔 & 표면파 전파", "METHOD")
card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(2.3),
     "시편 A (test #5-03)",
     ["", "· 5 µm 씩 이동하며 10회 측정",
      "· 손상부를 피해 라인 스캔", ""], BLUE, ts=16, bs=14)
card(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(2.3),
     "시편 B (test #5-23)",
     ["", "· 3 µm 씩 이동하며 5회 측정 (반복)",
      "· 파워 테스트 중 중앙부 손상 발생", ""], ORANGE, ts=16, bs=14)
bullets(s, [
    (0, "표면파 전파거리 = 펌프–프로브 간격 3 mm (고정)"),
    (0, "표면파 속도 = 전파거리 / 도달시간(ToF)  →  취득 신호에서 ToF 추정"),
    (1, "보유 데이터 기준 표면파 도달은 약 2.0–2.5 µs 구간에서 관측됨"),
    (0, "분석 데이터: 시편별 11개 위치 라인스캔, 파일당 100만 점 @ 5 GHz"),
], Inches(0.7), Inches(4.1), Inches(12), Inches(2.5), size=16, gap=10)
footer(s)

# ======================================================================
# Slide 7 — 결과 ① 취득 신호 (표면파 A-scan)
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "실험 결과 ① — 취득된 표면파 신호", "RESULT")
pic_fit(s, os.path.join(RESDIR, "ms_fig1_surfacewave.png"),
        Inches(0.55), Inches(1.4), Inches(8.1), Inches(5.2))
caption(s, Inches(0.55), Inches(6.55), Inches(8.1),
        "0–3 µs 윈도우 · 회색=개별위치, 진한선=평균, 빨강=포락선")
card(s, Inches(8.9), Inches(1.5), Inches(3.85), Inches(4.9),
     "읽는 법",
     ["", "· ~1 µs 부근부터 표면파", "  버스트가 도달", "",
      "· 포락선 피크 시점이", "  표면파 도달시간(ToF)", "",
      "· 두 시편 모두 표면파는", "  확인되나 진폭이 작고", "  노이즈가 상대적으로 큼",
      "", "→ SNR 개선 여지 큼"], BLUE, ts=15, bs=12)
footer(s)

# ======================================================================
# Slide 8 — 결과 ② 위치별 재현성 (waterfall)
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "실험 결과 ② — 위치별 신호 재현성", "RESULT")
pic_fit(s, os.path.join(RESDIR, "ms_fig2_waterfall.png"),
        Inches(0.55), Inches(1.4), Inches(8.1), Inches(5.2))
caption(s, Inches(0.55), Inches(6.55), Inches(8.1),
        "각 시편 11개 스캔 위치의 표면파 신호 적층")
card(s, Inches(8.9), Inches(1.5), Inches(3.85), Inches(4.9),
     "관찰",
     ["", "· 위치마다 파형이 흔들림", "  (위치 간 편차 존재)", "",
      "· 일부 위치는 표면파가", "  뚜렷, 일부는 미약", "",
      "· 원본 보고와 일치:", "  '일부 위치에서만 신호", "   명확히 관찰'", "",
      "→ 표면 상태/정렬 영향"], ORANGE, ts=15, bs=12)
footer(s)

# ======================================================================
# Slide 9 — 결과 ③ 스펙트럼 & 특징 비교
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "실험 결과 ③ — 시편 간 정량 비교", "RESULT")
pic_fit(s, os.path.join(RESDIR, "ms_fig3_spectrum.png"),
        Inches(0.5), Inches(1.4), Inches(6.3), Inches(2.7))
pic_fit(s, os.path.join(RESDIR, "ms_fig4_features.png"),
        Inches(0.5), Inches(4.2), Inches(6.3), Inches(2.9))
card(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3),
     "정량 비교에서 읽히는 경향  (잠정)",
     ["",
      "· 고주파(>2 MHz) 성분이 시편 B(조대립)", "  에서 더 낮음",
      "  → '결정립↑ ⇒ 고주파 감쇠↑' 가설과 부합",
      "",
      "· 신호 에너지도 시편 B가 다소 낮음",
      "",
      "· 도달시간/피크진폭/중심주파수는", "  두 시편 차이가 오차범위와 겹침",
      "",
      "⚠  주의: 차이가 작고 위치별 산란이 커",
      "    현 단계에서 결론으로 단정하기 어려움.",
      "    SNR 개선 후 재측정 필요."], TEAL, ts=15, bs=12)
footer(s)

# ======================================================================
# Slide 10 — 결과 분석 (신호 품질 이슈)
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "결과 분석 — 신호 품질 저하 요인", "ANALYSIS")
card(s, Inches(0.6), Inches(1.5), Inches(4.0), Inches(2.5),
     "1. 표면 거칠기",
     ["폴리싱 잔존 요철에서 레이저", "산란 발생 → 신호 취득 불안정,",
      "노이즈 증가로 판단됨."], BLUE, ts=16, bs=13)
card(s, Inches(4.75), Inches(1.5), Inches(4.0), Inches(2.5),
     "2. 시편 정렬 불량",
     ["고정 마운트의 수평/수직 정렬", "미흡 → 광축 어긋남으로",
      "신호 취득 불안정·노이즈 증가."], ORANGE, ts=16, bs=13)
card(s, Inches(8.9), Inches(1.5), Inches(4.0), Inches(2.5),
     "3. 레이저 출력 부족",
     ["시편 손상 방지를 위해 최소", "출력으로 발생 → 신호 자체가",
      "작아 SNR 저하."], TEAL, ts=16, bs=13)
rect(s, Inches(0.6), Inches(4.4), Inches(12.2), Inches(2.0), LIGHT)
tb, tf = textbox(s, Inches(0.9), Inches(4.6), Inches(11.6), Inches(1.7))
_set(tf.paragraphs[0], "종합", 17, NAVY, bold=True)
for ln in [
    "· 세 요인이 겹쳐 SNR이 낮아, 표면파는 검출되나 시편 간 미세조직 차이를",
    "  안정적으로 분리하기에는 신호 품질이 부족.",
    "· 결정립–표면파 상관성의 '경향'은 보이나(고주파 감쇠), 정량 결론에는 추가 측정 필요."]:
    p = tf.add_paragraph(); _set(p, ln, 14, GRAY); p.line_spacing = 1.15; p.space_before = Pt(3)
footer(s)

# ======================================================================
# Slide 11 — 추후 연구 방향
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "추후 연구 방향", "NEXT STEPS")
card(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(1.35),
     "1. 표면 거칠기 개선",
     ["연마재 입도를 낮춰(3 µm 수준) 추가 폴리싱 → 표면 조도 개선으로 산란·노이즈 저감"],
     BLUE, ts=16, bs=14)
card(s, Inches(0.6), Inches(3.0), Inches(12.2), Inches(1.35),
     "2. 정밀 정렬 도입",
     ["틸트/회전 자유도 포함 정밀 스테이지 + 카메라 기반 광축 정렬 확인으로 취득 안정화"],
     ORANGE, ts=16, bs=14)
card(s, Inches(0.6), Inches(4.5), Inches(12.2), Inches(1.35),
     "3. 출력 최적화",
     ["사전에 동일 시편으로 손상 임계 플루언스를 실험적으로 파악 → 출력을 점진적으로 증가시켜 SNR 확보"],
     TEAL, ts=16, bs=14)
tb, tf = textbox(s, Inches(0.6), Inches(6.05), Inches(12.2), Inches(0.7))
_set(tf.paragraphs[0],
     "+ 분석 측면: 결정립–표면파 속도/감쇠 상관모델 수립, 시편↔데이터 매핑 및 결정립 측정값 정합 확인",
     14, NAVY, bold=True)
footer(s)

# ======================================================================
# Slide 12 — 요약
# ======================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(1.7), SW, Pt(3), BLUE)
tb, tf = textbox(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(1.1))
_set(tf.paragraphs[0], "요약", 34, WHITE, bold=True)
pts = [
    "목표: 레이저 초음파(pump–probe SAW)로 결정립(64.4 vs 482 µm) 미세조직을 비파괴 분석",
    "결과: 표면파 신호는 검출되나 SNR이 낮고 위치별 편차 큼",
    "경향: 조대립 시편에서 고주파 성분·에너지가 다소 낮음 (가설과 부합, 단 잠정적)",
    "다음: 표면 폴리싱·정밀 정렬·출력 최적화로 SNR 확보 후 재측정 → 정량 상관 분석",
]
tb, tf = textbox(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(4))
for i, t in enumerate(pts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    _set(p, "•  " + t, 18, RGBColor(0xDD,0xE6,0xF2)); p.space_after = Pt(14); p.line_spacing = 1.15
tb, tf = textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7))
_set(tf.paragraphs[0], "감사합니다 — 질문 환영합니다 :)", 18, ACCENT2, bold=True)

out = os.path.join(ROOT, "기계연구원_샘플실험_재구성초안.pptx")
prs.save(out)
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
