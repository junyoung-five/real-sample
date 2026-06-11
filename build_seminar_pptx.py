# -*- coding: utf-8 -*-
"""
세미나용 PPTX 생성 스크립트
주제: VS Code + Git/GitHub + Claude Code 로 연구하기
사례: 레이저 초음파 결함검출 데이터 분석 (실제 진행 과정)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------- 테마 ----------------
NAVY   = RGBColor(0x1F, 0x2A, 0x44)
BLUE   = RGBColor(0x2E, 0x86, 0xDE)
TEAL   = RGBColor(0x12, 0xA5, 0x9B)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF9)
GRAY   = RGBColor(0x55, 0x5A, 0x66)
DGRAY  = RGBColor(0x33, 0x37, 0x40)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG = RGBColor(0x1E, 0x1E, 0x1E)
CODEFG = RGBColor(0xE6, 0xE6, 0xE6)
ACCENT2= RGBColor(0xF5, 0xA6, 0x23)

FONT   = "맑은 고딕"
MONO   = "Consolas"

RESDIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "analysis_matlab", "results")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(tf_para, text, size, color, bold=False, font=FONT, italic=False):
    tf_para.text = text
    for r in tf_para.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic


def rect(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def header(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Pt(4), BLUE)
    tb, tf = textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.95),
                     anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        p0 = tf.paragraphs[0]; _set(p0, kicker, 12, BLUE, bold=True)
        p1 = tf.add_paragraph(); _set(p1, title, 26, WHITE, bold=True)
    else:
        _set(tf.paragraphs[0], title, 28, WHITE, bold=True)


def footer(slide, n):
    tb, tf = textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    p = tf.paragraphs[0]
    _set(p, "Claude Code 활용 세미나   ·   레이저 초음파 데이터 분석 사례", 9, GRAY)
    tb2, tf2 = textbox(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    _set(p2, str(n), 10, GRAY, bold=True)


def bullets(slide, items, x, y, w, h, size=16, gap=10):
    tb, tf = textbox(slide, x, y, w, h)
    for i, (lvl, txt, *col) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        c = col[0] if col else (DGRAY if lvl == 0 else GRAY)
        bullet = "●  " if lvl == 0 else "–  "
        _set(p, bullet + txt, size if lvl == 0 else size-2, c, bold=(lvl == 0))
        p.level = lvl
        p.space_after = Pt(gap if lvl == 0 else gap-4)
        p.line_spacing = 1.1
    return tb


def codebox(slide, lines, x, y, w, h, size=12):
    sp = rect(slide, x, y, w, h, CODEBG)
    tb, tf = textbox(slide, x + Inches(0.18), y + Inches(0.12),
                     w - Inches(0.36), h - Inches(0.24))
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        color = CODEFG
        if ln.strip().startswith("#") or ln.strip().startswith("%"):
            color = RGBColor(0x6A, 0x99, 0x55)      # comment green
        elif ln.strip().startswith(">") or ln.strip().startswith("$"):
            color = RGBColor(0x9C, 0xDC, 0xFE)      # prompt blue
        _set(p, ln, size, color, font=MONO)
        p.line_spacing = 1.05
    return sp


def card(slide, x, y, w, h, title, body, accent=BLUE, ts=15, bs=12):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, Inches(0.12), h, accent)
    tb, tf = textbox(slide, x + Inches(0.28), y + Inches(0.12),
                     w - Inches(0.45), h - Inches(0.2))
    _set(tf.paragraphs[0], title, ts, NAVY, bold=True)
    for line in body:
        p = tf.add_paragraph(); _set(p, line, bs, GRAY); p.line_spacing = 1.08
        p.space_before = Pt(3)


# ======================================================================
# Slide 1 — 표지
# ======================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Pt(3), BLUE)
tb, tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(3))
_set(tf.paragraphs[0], "AI 와 함께 연구하기", 20, BLUE, bold=True)
p = tf.add_paragraph(); _set(p, "VS Code · Git/GitHub · Claude Code", 44, WHITE, bold=True)
p.space_before = Pt(6)
p = tf.add_paragraph()
_set(p, "코드를 직접 짜지 않고, 대화로 데이터 분석을 끝내기", 18, RGBColor(0xC9,0xD6,0xEA))
p.space_before = Pt(14)
tb, tf = textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(2))
_set(tf.paragraphs[0], "실제 사례 :  레이저 초음파 시편 결함검출 — 오실로스코프 원본 데이터 분석",
     15, ACCENT2, bold=True)
p = tf.add_paragraph(); _set(p, "연구실 세미나   ·   10–15분   ·   입문자 대상", 13,
                             RGBColor(0xAE,0xB9,0xCC)); p.space_before = Pt(8)

# ======================================================================
# Slide 2 — 오늘의 목표
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "오늘 세미나의 목표", "WHY")
bullets(s, [
    (0, "코딩이 전공이 아니어도, 연구에 필요한 분석 코드를 빠르게 만들 수 있다"),
    (0, "세 가지 도구가 어떻게 맞물려 돌아가는지 '큰 그림'을 잡는다"),
    (1, "VS Code = 작업 공간 ·  Git/GitHub = 기록·공유 ·  Claude Code = AI 조수"),
    (0, "실제 레이저 초음파 데이터 분석이 어떻게 진행됐는지 직접 본다"),
    (0, "내 연구에 오늘 바로 적용할 수 있는 출발점을 얻는다"),
], Inches(0.7), Inches(1.6), Inches(12), Inches(5), size=19, gap=16)
footer(s, 2)

# ======================================================================
# Slide 3 — 발표 순서
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "발표 순서", "AGENDA")
items = [
    ("1", "왜 필요한가", "연구자의 코딩 현실"),
    ("2", "도구 3종 소개", "VS Code · Git/GitHub · Claude Code"),
    ("3", "동작 원리", "대화 → 코드 → 실행 → 결과"),
    ("4", "실제 사례", "레이저 초음파 데이터 분석 전 과정"),
    ("5", "실전 팁 & 주의", "좋은 프롬프트 · 검증 · 시작하기"),
]
y = Inches(1.55)
for num, t, d in items:
    rect(s, Inches(0.7), y, Inches(0.7), Inches(0.7), BLUE)
    tb, tf = textbox(s, Inches(0.7), y, Inches(0.7), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER; _set(pp, num, 20, WHITE, bold=True)
    tb, tf = textbox(s, Inches(1.65), y, Inches(10.8), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], t + "   —   " + d, 18, DGRAY, bold=True)
    y += Inches(1.02)
footer(s, 3)

# ======================================================================
# Slide 4 — 연구자의 코딩 현실 (문제 제기)
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "연구자의 코딩 현실", "WHY")
card(s, Inches(0.6), Inches(1.5), Inches(5.9), Inches(1.65),
     "“분석 코드 짜는 데 시간이 다 간다”",
     ["측정 데이터 포맷 파악, 필터·FFT·시각화…", "본 연구보다 코드에 시간이 더 든다"], BLUE)
card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(1.65),
     "“새 언어·툴은 배울 엄두가 안 난다”",
     ["MATLAB? Python? 문법부터 막힌다", "예제 찾고 검색하다 하루가 간다"], TEAL)
card(s, Inches(0.6), Inches(3.4), Inches(5.9), Inches(1.65),
     "“에러 한 줄에 반나절”",
     ["NaN, 차원 안 맞음, 라이브러리 충돌…", "디버깅이 가장 큰 병목"], ACCENT2)
card(s, Inches(6.8), Inches(3.4), Inches(5.9), Inches(1.65),
     "“재현·인수인계가 안 된다”",
     ["내 PC에서만 돌아가는 코드", "졸업하면 사라지는 분석 노하우"], NAVY)
rect(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.0), LIGHT)
tb, tf = textbox(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.0),
                 anchor=MSO_ANCHOR.MIDDLE)
_set(tf.paragraphs[0],
     "→  이 네 가지를 VS Code · Git/GitHub · Claude Code 조합이 정면으로 해결한다",
     18, NAVY, bold=True)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
footer(s, 4)

# ======================================================================
# Slide 5 — 도구 3종 한눈에
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "세 가지 도구, 역할 분담", "OVERVIEW")
card(s, Inches(0.6), Inches(1.6), Inches(3.95), Inches(4.6),
     "VS Code", ["",
     "코드 작성·실행이 이뤄지는", "작업 공간(에디터)", "",
     "· 무료, 가볍고 빠름", "· 통합 터미널 내장", "· 확장으로 기능 추가",
     "", "→ 모든 작업의 '책상'"], BLUE, ts=22, bs=13)
card(s, Inches(4.7), Inches(1.6), Inches(3.95), Inches(4.6),
     "Git / GitHub", ["",
     "코드의 변경 이력을", "기록·백업·공유", "",
     "· 언제든 과거로 복구", "· 클라우드 백업", "· 동료와 협업·재현",
     "", "→ 연구의 '실험노트'"], TEAL, ts=22, bs=13)
card(s, Inches(8.8), Inches(1.6), Inches(3.95), Inches(4.6),
     "Claude Code", ["",
     "대화로 코드를 짜고 실행해", "주는 AI 에이전트", "",
     "· 파일 읽기·쓰기·수정", "· 직접 실행·디버깅", "· 자연어로 지시",
     "", "→ 옆자리 '코딩 조수'"], ACCENT2, ts=22, bs=13)
footer(s, 5)

# ======================================================================
# Slide 6 — VS Code
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "VS Code 란?", "도구 ①")
bullets(s, [
    (0, "Microsoft가 만든 무료 코드 에디터 (가장 널리 쓰임)"),
    (0, "한 화면에서 코드 작성·실행·결과 확인을 모두 처리"),
    (1, "왼쪽: 파일 탐색기   /   가운데: 코드   /   아래: 통합 터미널"),
    (0, "'확장(Extension)'으로 원하는 언어·기능을 자유롭게 추가"),
    (1, "Python, MATLAB, Git, 그리고 Claude Code 까지 한 곳에"),
    (0, "오늘 본 분석도 전부 VS Code 안에서 이뤄짐 (창을 떠나지 않음)"),
], Inches(0.7), Inches(1.55), Inches(7.3), Inches(5), size=17, gap=14)
card(s, Inches(8.3), Inches(1.6), Inches(4.4), Inches(4.4),
     "핵심 한 줄",
     ["", "연구에 필요한 모든 작업이", "모이는 '디지털 작업 책상'.",
      "", "에디터 + 터미널 + 확장이", "하나로 통합되어 있어",
      "도구를 옮겨다닐 필요가 없다."], BLUE, ts=18, bs=14)
footer(s, 6)

# ======================================================================
# Slide 7 — Git/GitHub
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "Git / GitHub 란?", "도구 ②")
bullets(s, [
    (0, "Git : 코드의 '변경 이력'을 스냅샷으로 저장하는 버전관리 도구"),
    (1, "잘못돼도 언제든 과거 시점으로 되돌릴 수 있다"),
    (0, "GitHub : 그 이력을 올려두는 클라우드 저장소 (백업 + 공유)"),
    (0, "연구에서 특히 중요한 이유"),
    (1, "재현성: '어떤 코드로 이 결과를 냈는가'가 기록으로 남음"),
    (1, "협업: 동료와 같은 코드를 안전하게 함께 수정"),
    (1, "인수인계: 졸업·이직 후에도 분석이 사라지지 않음"),
], Inches(0.7), Inches(1.55), Inches(7.3), Inches(5), size=17, gap=12)
card(s, Inches(8.3), Inches(1.6), Inches(4.4), Inches(4.4),
     "비유",
     ["", "Git = 게임의 '세이브 파일'", "",
      "위험한 시도 전에 저장하고,", "실패하면 불러오기.", "",
      "GitHub = 그 세이브를", "클라우드에 올려 동료와", "공유하는 것."], TEAL, ts=18, bs=14)
footer(s, 7)

# ======================================================================
# Slide 8 — Claude Code
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "Claude Code 란?", "도구 ③")
bullets(s, [
    (0, "터미널/VS Code 안에서 동작하는 AI 코딩 에이전트"),
    (0, "일반 챗봇과 결정적 차이 : 말만 하는 게 아니라 '직접 한다'"),
    (1, "내 폴더의 파일을 직접 열어 읽고 (데이터 포맷 파악)"),
    (1, "코드 파일을 직접 만들고 수정하고"),
    (1, "프로그램을 직접 실행해 결과를 확인하고 버그를 고친다"),
    (0, "자연어로 지시 → 사람은 '무엇을' 원하는지만 말하면 됨"),
], Inches(0.7), Inches(1.55), Inches(7.3), Inches(5), size=17, gap=13)
card(s, Inches(8.3), Inches(1.6), Inches(4.4), Inches(4.4),
     "오늘의 놀라운 장면",
     ["", "MATLAB GUI 를 한 번도", "띄우지 않고,", "",
      "Claude 가 matlab 을", "직접 실행(-batch)해", "결과·그래프까지 뽑고",
      "버그도 스스로 고침."], ACCENT2, ts=18, bs=14)
footer(s, 8)

# ======================================================================
# Slide 9 — 워크플로우
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "동작 원리 : 대화 한 번의 사이클", "HOW IT WORKS")
steps = [
    ("①", "요청", "자연어로 지시\n“이 데이터 분석\n코드 짜줘”", BLUE),
    ("②", "탐색·작성", "파일을 직접 읽고\n코드를 생성", TEAL),
    ("③", "실행·검증", "직접 실행하고\n에러를 디버깅", ACCENT2),
    ("④", "결과·반복", "그래프·수치 확인\n필요시 다시 수정", NAVY),
]
x = Inches(0.55); w = Inches(2.85); y = Inches(2.1); h = Inches(2.7)
for i, (num, t, d, c) in enumerate(steps):
    rect(s, x, y, w, h, LIGHT)
    rect(s, x, y, w, Inches(0.7), c)
    tb, tf = textbox(s, x, y, w, Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    _set(pp, num + "  " + t, 18, WHITE, bold=True)
    tb, tf = textbox(s, x + Inches(0.15), y + Inches(0.85), w - Inches(0.3), Inches(1.7),
                     anchor=MSO_ANCHOR.TOP)
    for j, line in enumerate(d.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        _set(p, line, 14, GRAY); p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.1
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                x + w + Inches(0.02), y + Inches(1.05),
                                Inches(0.32), Inches(0.5))
        ar.fill.solid(); ar.fill.fore_color.rgb = BLUE; ar.line.fill.background()
        ar.shadow.inherit = False
    x = x + w + Inches(0.36)
rect(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.0), NAVY)
tb, tf = textbox(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.0),
                 anchor=MSO_ANCHOR.MIDDLE)
_set(tf.paragraphs[0],
     "사람은 '무엇을·왜'에 집중하고, '어떻게(코드·실행·디버깅)'는 AI 가 담당",
     18, WHITE, bold=True); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
footer(s, 9)

# ======================================================================
# Slide 10 — 사례 소개
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "실제 사례 : 레이저 초음파 결함검출", "CASE STUDY")
bullets(s, [
    (0, "목표 : 레이저 초음파로 시편 내부 결함을 검출하는 실험 데이터 분석"),
    (0, "가진 것 : 오실로스코프(LeCroy)에서 추출한 원본 .dat 파일 22개"),
    (1, "시편 2종 × 라인스캔 11위치"),
    (1, "파일당 100만 포인트, 시간 −20~180 µs"),
    (0, "출발점에 내가 아는 것 : 거의 없음 (포맷·샘플링조차 직접 안 봄)"),
    (0, "한 일 : Claude 에게 폴더를 주고 \"분석 코드 짜줘\" 라고 요청"),
], Inches(0.7), Inches(1.55), Inches(7.4), Inches(5), size=17, gap=13)
card(s, Inches(8.4), Inches(1.6), Inches(4.3), Inches(4.4),
     "데이터 한눈에",
     ["", "·  포맷: 2열 ASCII (시간, 전압)", "·  파일: C1--sample#--NNNNN.dat",
      "·  포인트: 1,000,002 개/파일", "·  샘플링: 5 GHz (Δt=0.2 ns)",
      "·  시편: sample1, sample2", "·  위치: 00000 ~ 00010"], BLUE, ts=17, bs=13)
footer(s, 10)

# ======================================================================
# Slide 11 — Step 1: 요청
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "Step 1 — 자연어로 요청하기", "CASE STUDY")
tb, tf = textbox(s, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.5))
_set(tf.paragraphs[0], "전문 용어도, 코드도 필요 없다. 평소 말하듯 요청한다.", 16, GRAY)
codebox(s, [
    "> 현재 열려있는 프로젝트 경로는 레이저 초음파를 활용해 시편의",
    "  결함을 검출하는 실험에서 오실로스코프로 추출한 원본 데이터다.",
    "",
    "> 이 데이터들을 분석하기 위한 방법론들을 MATLAB 코드로",
    "  구현해서 데이터 분석을 진행하고 싶어.",
    "",
    "> 데이터를 참고해서 적절한 분석용 MATLAB 코드를 작성해줘.",
], Inches(0.7), Inches(2.1), Inches(11.9), Inches(2.7), size=15)
card(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.35),
     "포인트",
     ["폴더 경로만 알려주면 Claude 가 알아서 파일을 열어 본다. ",
      "'어떻게 짜라'가 아니라 '무엇을 하고 싶다'만 말하면 된다."], TEAL, ts=15, bs=14)
footer(s, 11)

# ======================================================================
# Slide 12 — Step 2: 스스로 파악 + 되물음
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "Step 2 — AI 가 스스로 파악하고, 모르면 되묻는다", "CASE STUDY")
card(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(2.2),
     "✔  스스로 알아낸 것",
     ["· 파일 22개 / 2열 ASCII 포맷 자동 인식",
      "· 샘플링 5 GHz, 100만 포인트 계산",
      "· t=0 = 레이저 발사 시점(트리거) 추론",
      "· t<0 구간은 노이즈 → SNR 추정에 활용"], TEAL, ts=16, bs=13)
card(s, Inches(6.8), Inches(1.55), Inches(5.9), Inches(2.2),
     "❓  애매해서 사람에게 물어본 것",
     ["· 11개 파일 = 반복측정? 위치별 스캔?",
      "· 두 시편의 관계는? (기준 vs 결함?)",
      "", "→ 객관식으로 질문 →  답: 라인스캔 /",
      "    자동비교  →  분석 방향 확정"], ACCENT2, ts=16, bs=13)
tb, tf = textbox(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.5))
_set(tf.paragraphs[0], "실제로 받은 질문 화면 :", 14, GRAY, bold=True)
codebox(s, [
    "  ? 한 시편당 11개 파일(00000~00010)은 어떤 성격인가요?",
    "      ○ 동일 지점 반복 측정      ● 라인 스캔(위치별)",
    "  ? 두 시편(sample1, sample2)의 관계는?",
    "      ● 잘 모름 / 자동 비교      ○ 기준(건전) vs 결함",
], Inches(0.6), Inches(4.5), Inches(12.1), Inches(1.7), size=14)
footer(s, 12)

# ======================================================================
# Slide 13 — Step 3: 코드 생성
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "Step 3 — 모듈식 분석 코드 자동 생성", "CASE STUDY")
tb, tf = textbox(s, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.5))
_set(tf.paragraphs[0], "한 덩어리가 아니라, 역할별로 나뉜 9개 파일을 체계적으로 생성", 16, GRAY)
files = [
    ("lu_config.m", "모든 파라미터 중앙 설정"),
    ("lu_load_dat.m", "오실로스코프 파일 읽기"),
    ("lu_preprocess.m", "노이즈 제거·필터·SNR"),
    ("lu_time_analysis.m", "도달시간(ToF)·에코·에너지"),
    ("lu_freq_analysis.m", "FFT·스펙트로그램"),
    ("lu_bscan.m", "B-scan 영상 구성"),
    ("lu_compare.m", "결함 후보 자동 판정"),
    ("lu_plots.m", "결과 시각화"),
    ("run_analysis.m", "전체 실행 (메인)"),
]
x0, y0 = Inches(0.7), Inches(2.1); cw, ch = Inches(3.9), Inches(0.92)
for i, (fn, desc) in enumerate(files):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + Inches(0.18)); y = y0 + r * (ch + Inches(0.16))
    rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, Inches(0.1), ch, BLUE)
    tb, tf = textbox(s, x + Inches(0.22), y + Inches(0.08), cw - Inches(0.3), ch - Inches(0.1))
    _set(tf.paragraphs[0], fn, 14, NAVY, bold=True, font=MONO)
    p = tf.add_paragraph(); _set(p, desc, 12, GRAY)
tb, tf = textbox(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.6))
_set(tf.paragraphs[0],
     "→ 모듈식 구조라 나중에 한 부분만 수정·재사용하기 쉽다 (좋은 코드 습관까지 제공)",
     14, TEAL, bold=True)
footer(s, 13)

# ======================================================================
# Slide 14 — Step 4: 직접 실행 + 디버깅 (압권)
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "Step 4 — 직접 실행하고, 버그까지 스스로 고친다", "CASE STUDY ★")
tb, tf = textbox(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.5))
_set(tf.paragraphs[0], "GUI 없이 터미널에서 MATLAB 을 직접 호출 → 실행 → 문제 발견 → 수정 → 재실행", 15, GRAY)
codebox(s, [
    "$ matlab -batch \"run_analysis\"      % GUI 없이 직접 실행",
    "  ...",
    "  % 결과: 모든 값이 NaN  ←  Claude 가 스스로 이상 감지",
], Inches(0.7), Inches(1.95), Inches(11.9), Inches(1.35), size=13)
card(s, Inches(0.7), Inches(3.5), Inches(5.9), Inches(2.45),
     "발견한 버그",
     ["5 GHz 에서 0.5 MHz 필터는",
      "정규화 주파수가 0.0002 로 너무 낮아",
      "필터가 수치적으로 불안정 → NaN.",
      "", "(사람이면 한참 헤맬 함정)"], ACCENT2, ts=16, bs=13)
card(s, Inches(6.8), Inches(3.5), Inches(5.9), Inches(2.45),
     "스스로 한 해결",
     ["① 필터를 SOS(2차 섹션)로 재설계",
      "   → 수치 안정화",
      "② 신호대역에 맞춰 데시메이션 적용",
      "③ 재실행 → 정상 결과 확인", ""], TEAL, ts=16, bs=13)
footer(s, 14)

# ======================================================================
# Slide 15 — Step 5: 결과물 (이미지)
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "Step 5 — 결과물 : 그래프와 결함 후보까지", "CASE STUDY")
bscan = os.path.join(RESDIR, "fig4_bscan.png")
feat  = os.path.join(RESDIR, "fig5_features.png")
if os.path.exists(bscan):
    s.shapes.add_picture(bscan, Inches(0.55), Inches(1.5), height=Inches(3.0))
if os.path.exists(feat):
    s.shapes.add_picture(feat, Inches(0.55), Inches(4.55), height=Inches(2.7))
card(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.25),
     "자동으로 얻은 것",
     ["",
      "·  B-scan 영상 (위치 × 시간 에너지)",
      "·  A-scan·포락선·FFT·스펙트로그램",
      "·  위치별 특징곡선 (진폭·ToF·주파수)",
      "",
      "·  결함 후보 자동 검출 (통계적 이상치)",
      "    → sample1 위치 8 이 이상치로 표시",
      "·  주파수 1.0→0.6 MHz 하락 등",
      "    해석 가능한 물리적 패턴 포착",
      "",
      "·  그림 5종(PNG) · CSV · MAT 파일 저장",
      "",
      "→ 파라미터(두께·음속·간격)만 넣으면",
      "   정량적 깊이까지 바로 확장 가능"], BLUE, ts=16, bs=13)
footer(s, 15)

# ======================================================================
# Slide 16 — Git/GitHub로 마무리
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "마무리 : Git/GitHub 로 기록하고 공유", "WORKFLOW")
tb, tf = textbox(s, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.5))
_set(tf.paragraphs[0], "생성된 코드와 결과를 버전관리에 올리면 '재현 가능한 연구'가 완성된다.", 16, GRAY)
codebox(s, [
    "$ git init                       # 저장소 시작",
    "$ git add .                      # 변경사항 담기",
    "$ git commit -m \"LU 분석 코드 추가\"  # 스냅샷 저장",
    "$ git push                       # GitHub 에 백업·공유",
    "",
    "# 이 모든 명령도 Claude 에게 \"커밋하고 푸시해줘\" 라고",
    "# 부탁하면 대신 실행해 준다.",
], Inches(0.7), Inches(2.05), Inches(11.9), Inches(2.7), size=14)
card(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.45),
     "얻는 것",
     ["언제든 과거로 복구 · 동료와 협업 · 졸업 후에도 남는 분석 · ",
      "'이 코드로 이 결과를 냈다'는 완전한 재현 기록"], TEAL, ts=15, bs=14)
footer(s, 16)

# ======================================================================
# Slide 17 — 좋은 프롬프트 팁
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "실전 팁 : 좋은 요청 = 좋은 결과", "TIPS")
card(s, Inches(0.6), Inches(1.55), Inches(5.95), Inches(2.0),
     "① 맥락을 준다",
     ["무엇을 하는 데이터인지, 목표가 무엇인지", "한두 줄이라도 배경을 설명한다."], BLUE)
card(s, Inches(6.78), Inches(1.55), Inches(5.95), Inches(2.0),
     "② 구체적으로 말한다",
     ["\"분석해줘\" 보다 \"FFT 로 주파수 분석하고", "그래프로 보여줘\" 가 훨씬 낫다."], TEAL)
card(s, Inches(0.6), Inches(3.7), Inches(5.95), Inches(2.0),
     "③ 단계적으로 진행한다",
     ["한 번에 다 시키기보다, 만들고→확인하고", "→ 다음을 요청하며 키워간다."], ACCENT2)
card(s, Inches(6.78), Inches(3.7), Inches(5.95), Inches(2.0),
     "④ 검증을 요청한다",
     ["\"실제로 실행해서 결과 확인해줘\",", "\"테스트도 만들어줘\" 라고 명시한다."], NAVY)
rect(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.65), LIGHT)
tb, tf = textbox(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.65),
                 anchor=MSO_ANCHOR.MIDDLE)
_set(tf.paragraphs[0], "모르면 \"이거 왜 이렇게 했어?\" 라고 되물어도 된다 — 설명까지 해준다.",
     15, NAVY, bold=True); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
footer(s, 17)

# ======================================================================
# Slide 18 — 한계와 주의
# ======================================================================
s = prs.slides.add_slide(BLANK); header(s, "꼭 기억할 한계와 주의점", "CAUTION")
bullets(s, [
    (0, "결과 검증은 사람의 몫 — AI도 틀린다. 그래프·수치를 직접 확인할 것"),
    (0, "도메인 지식은 연구자가 — 시편 두께·음속·스캔 간격 같은 실험 파라미터는"),
    (1, "사람이 넣어야 정량 분석이 유효 (오늘도 그 값이 없어 정성 분석까지만)"),
    (0, "민감/미공개 데이터는 주의 — 무엇을 올리고 공유하는지 항상 인지"),
    (0, "AI는 '가속기'지 '대체재'가 아니다 — 연구의 판단과 해석은 여전히 사람"),
], Inches(0.7), Inches(1.6), Inches(12), Inches(4), size=18, gap=15)
footer(s, 18)

# ======================================================================
# Slide 19 — 시작하기 / 마무리
# ======================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.0), SW, Pt(3), BLUE)
tb, tf = textbox(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.3))
_set(tf.paragraphs[0], "오늘부터 시작하기", 38, WHITE, bold=True)
steps = [
    ("1", "VS Code 설치", "code.visualstudio.com 에서 무료 다운로드"),
    ("2", "Claude Code 설치", "터미널에서 안내대로 설치 후 로그인"),
    ("3", "폴더 열고 대화 시작", "내 데이터 폴더를 열고 하고 싶은 걸 말해보기"),
]
y = Inches(2.4)
for num, t, d in steps:
    rect(s, Inches(0.9), y, Inches(0.65), Inches(0.65), BLUE)
    tb, tf = textbox(s, Inches(0.9), y, Inches(0.65), Inches(0.65), anchor=MSO_ANCHOR.MIDDLE)
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER; _set(pp, num, 18, WHITE, bold=True)
    tb, tf = textbox(s, Inches(1.75), y, Inches(10.8), Inches(0.65), anchor=MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], t, 19, WHITE, bold=True)
    p = tf.add_paragraph(); _set(p, d, 14, RGBColor(0xAE,0xB9,0xCC))
    y += Inches(0.95)
tb, tf = textbox(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.2))
_set(tf.paragraphs[0], "감사합니다 — 질문 환영합니다 :)", 24, ACCENT2, bold=True)
p = tf.add_paragraph()
_set(p, "“코드를 못 짜도, 분석은 시작할 수 있습니다.”", 16, RGBColor(0xC9,0xD6,0xEA))
p.space_before = Pt(8)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ClaudeCode_세미나.pptx")
prs.save(out)
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
